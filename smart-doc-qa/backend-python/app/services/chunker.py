"""
Document ingestion: parse → pre-process → smart chunk → return chunks with metadata.
Supports: PDF, DOCX, PPTX, XLSX, TXT, Images (OCR).
"""

import os
import uuid
from pathlib import Path
from typing import List, Dict, Any
from dataclasses import dataclass, field

import pytesseract
from PIL import Image


@dataclass
class Chunk:
    """A single text chunk with rich metadata."""
    chunk_id: str = field(default_factory=lambda: str(uuid.uuid4()))
    document_id: str = ""
    content: str = ""
    page_number: int = 0
    section_heading: str = ""
    chunk_index: int = 0
    total_chunks: int = 0
    document_name: str = ""
    file_type: str = ""
    # Hierarchical: every small chunk also stores its parent (larger) context
    parent_content: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "chunk_id": self.chunk_id,
            "document_id": self.document_id,
            "content": self.content,
            "page_number": self.page_number,
            "section_heading": self.section_heading,
            "chunk_index": self.chunk_index,
            "total_chunks": self.total_chunks,
            "document_name": self.document_name,
            "file_type": self.file_type,
            "parent_content": self.parent_content,
        }


# ── Text extraction per file type ─────────────────────────────────────────────

def _extract_pdf(path: str) -> List[Dict]:
    """Returns list of {page_number, text, heading}."""
    import PyPDF2
    pages = []
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if not text.strip():
                # Scanned PDF – fall back to OCR via image rendering
                # Requires pdf2image: pip install pdf2image
                try:
                    from pdf2image import convert_from_path
                    imgs = convert_from_path(path, first_page=i + 1, last_page=i + 1)
                    text = pytesseract.image_to_string(imgs[0]) if imgs else ""
                except Exception:
                    pass
            pages.append({"page_number": i + 1, "text": text, "heading": ""})
    return pages


def _extract_docx(path: str) -> List[Dict]:
    from docx import Document
    doc = Document(path)
    pages, current_page, current_heading = [], [], ""
    page_num = 1
    for para in doc.paragraphs:
        style = para.style.name or ""
        if style.startswith("Heading"):
            current_heading = para.text.strip()
        current_page.append(para.text)
        # Simulate page break every ~50 paragraphs
        if len(current_page) >= 50:
            pages.append({"page_number": page_num, "text": "\n".join(current_page), "heading": current_heading})
            current_page = []
            page_num += 1
    if current_page:
        pages.append({"page_number": page_num, "text": "\n".join(current_page), "heading": current_heading})
    return pages


def _extract_pptx(path: str) -> List[Dict]:
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        heading = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.shape_type == 13:  # title
                    heading = shape.text.strip()
                else:
                    texts.append(shape.text.strip())
        pages.append({"page_number": i + 1, "text": "\n".join(texts), "heading": heading})
    return pages


def _extract_xlsx(path: str) -> List[Dict]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    for sheet_idx, sheet in enumerate(wb.worksheets):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append("\t".join(str(c) if c is not None else "" for c in row))
        pages.append({
            "page_number": sheet_idx + 1,
            "text": "\n".join(rows),
            "heading": sheet.title,
        })
    return pages


def _extract_txt(path: str) -> List[Dict]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [{"page_number": 1, "text": text, "heading": ""}]


def _extract_image(path: str) -> List[Dict]:
    img = Image.open(path)
    text = pytesseract.image_to_string(img)
    return [{"page_number": 1, "text": text, "heading": ""}]


EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".txt":  _extract_txt,
    ".png":  _extract_image,
    ".jpg":  _extract_image,
    ".jpeg": _extract_image,
}


# ── Semantic / hierarchical chunking ──────────────────────────────────────────

def _split_into_sentences(text: str) -> List[str]:
    """Simple sentence splitter (replace with spaCy for production)."""
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    return [s.strip() for s in sentences if s.strip()]


def _hierarchical_chunk(
    page_text: str,
    small_size: int = 256,
    large_size: int = 512,
    overlap: int = 64,
) -> List[Dict[str, str]]:
    """
    Returns list of {small_chunk, large_chunk (parent context)}.
    small_chunk  → used for precise retrieval
    large_chunk  → passed to LLM for richer context
    """
    words = page_text.split()
    if not words:
        return []

    chunks = []
    # Build large chunks first
    large_chunks = []
    i = 0
    while i < len(words):
        large_chunks.append(" ".join(words[i: i + large_size]))
        i += large_size - overlap

    # Each large chunk → multiple small chunks
    for large in large_chunks:
        large_words = large.split()
        j = 0
        while j < len(large_words):
            small = " ".join(large_words[j: j + small_size])
            if small.strip():
                chunks.append({"small": small, "large": large})
            j += small_size - overlap

    return chunks


# ── Public API ────────────────────────────────────────────────────────────────

def ingest_document(
    file_path: str,
    document_id: str,
    document_name: str,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
) -> List[Chunk]:
    """
    Parse a document and return a flat list of Chunk objects ready for embedding.
    """
    ext = Path(file_path).suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")

    pages = extractor(file_path)
    all_chunks: List[Chunk] = []

    for page in pages:
        page_text = page["text"]
        heading = page.get("heading", "")
        page_num = page["page_number"]

        if not page_text.strip():
            continue

        raw_chunks = _hierarchical_chunk(page_text, chunk_size, chunk_size * 2, chunk_overlap)

        for idx, c in enumerate(raw_chunks):
            all_chunks.append(Chunk(
                document_id=document_id,
                content=c["small"],
                parent_content=c["large"],
                page_number=page_num,
                section_heading=heading,
                chunk_index=len(all_chunks),
                document_name=document_name,
                file_type=ext.lstrip("."),
            ))

    # Patch total_chunks
    total = len(all_chunks)
    for c in all_chunks:
        c.total_chunks = total

    return all_chunks
